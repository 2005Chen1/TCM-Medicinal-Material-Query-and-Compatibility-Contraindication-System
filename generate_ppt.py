"""
生成《中医药常用药材查询与配伍禁忌系统》需求分析汇报PPT
重构版：5章结构、数据库分页、架构图拆分、章节过渡页
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Constants ──────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(BASE, "附录-图")
OUTPUT = os.path.join(BASE, "需求分析汇报.pptx")

DEEP_RED = RGBColor(0x8B, 0x00, 0x00)
CREAM = RGBColor(0xFF, 0xF8, 0xF0)
GOLD = RGBColor(0xC9, 0xA9, 0x6E)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF2, 0xED)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_BG = RGBColor(0x3D, 0x1A, 0x1A)

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ────────────────────────────────────────────────────
def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def tbox(slide, left, top, width, height, text="", size=Pt(18),
         color=DARK, bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY,
         anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if p.runs:
        r = p.runs[0]
        r.font.size = size
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb


def rtf(slide, left, top, width, height):
    """Create a rich text frame for multi-paragraph use"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def add_p(tf, text, size=Pt(16), color=DARK, bold=False, align=PP_ALIGN.LEFT,
          font=FONT_BODY, before=Pt(3), after=Pt(3)):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = before
    p.space_after = after
    if p.runs:
        r = p.runs[0]
        r.font.size = size
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return p


def gold_line(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()


def section_title(slide, title, subtitle=""):
    """Standard content slide title with gold accent bar"""
    add_rect(slide, Inches(0.8), Inches(0.5), Pt(5), Inches(0.55), GOLD)
    tbox(slide, Inches(1.1), Inches(0.4), Inches(11), Inches(0.65),
         title, size=Pt(30), color=DEEP_RED, bold=True, font=FONT_TITLE)
    gold_line(slide, Inches(1.1), Inches(1.05), Inches(10.5))
    if subtitle:
        tbox(slide, Inches(1.1), Inches(1.15), Inches(10.5), Inches(0.35),
             subtitle, size=Pt(13), color=MID_GRAY)


def footer(slide, num):
    gold_line(slide, Inches(0.8), Inches(7.05), Inches(11.5))
    tbox(slide, Inches(0.8), Inches(7.1), Inches(6), Inches(0.3),
         "中医药常用药材查询与配伍禁忌系统 · 需求分析汇报",
         size=Pt(8), color=MID_GRAY)
    tbox(slide, Inches(12.1), Inches(7.1), Inches(1), Inches(0.3),
         str(num), size=Pt(10), color=MID_GRAY, align=PP_ALIGN.RIGHT)


def chapter_divider(chapter_num, cn_title, en_title, description):
    """Chapter transition slide - dark red bg, centered"""
    s = add_blank_slide()
    add_bg(s, DEEP_RED)
    add_rect(s, Inches(0), Inches(0), W, Pt(6), GOLD)
    add_rect(s, Inches(0), H - Pt(6), W, Pt(6), GOLD)

    # Chapter number
    tbox(s, Inches(3), Inches(1.6), Inches(7.333), Inches(1.0),
         f"第{chapter_num}章", size=Pt(28), color=GOLD, bold=True,
         font=FONT_TITLE, align=PP_ALIGN.CENTER)

    gold_line(s, Inches(5.2), Inches(2.6), Inches(2.933))

    # Chapter title
    tbox(s, Inches(2), Inches(2.9), Inches(9.333), Inches(1.0),
         cn_title, size=Pt(40), color=WHITE, bold=True,
         font=FONT_TITLE, align=PP_ALIGN.CENTER)

    # English title
    tbox(s, Inches(2), Inches(3.8), Inches(9.333), Inches(0.5),
         en_title, size=Pt(16), color=RGBColor(0xCC, 0xAA, 0x88),
         font=FONT_BODY, align=PP_ALIGN.CENTER)

    # Description of what's in this chapter
    tbox(s, Inches(3), Inches(4.7), Inches(7.333), Inches(0.8),
         description, size=Pt(14), color=RGBColor(0xDD, 0xCC, 0xBB),
         font=FONT_BODY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, DEEP_RED)
add_rect(s, Inches(0), Inches(0), W, Pt(8), GOLD)
add_rect(s, Inches(0), H - Pt(8), W, Pt(8), GOLD)
add_rect(s, Inches(2.0), Inches(1.8), Inches(9.333), Pt(1.5), GOLD)
add_rect(s, Inches(2.0), Inches(5.5), Inches(9.333), Pt(1.5), GOLD)

tbox(s, Inches(1.5), Inches(2.2), Inches(10.333), Inches(1.1),
     "中医药常用药材查询\n与配伍禁忌系统", size=Pt(44), color=WHITE,
     bold=True, font=FONT_TITLE, align=PP_ALIGN.CENTER)
gold_line(s, Inches(5), Inches(3.6), Inches(3.333))
tbox(s, Inches(1.5), Inches(3.9), Inches(10.333), Inches(0.7),
     "需 求 分 析 汇 报", size=Pt(28), color=GOLD, font=FONT_TITLE,
     align=PP_ALIGN.CENTER)
tbox(s, Inches(1.5), Inches(4.8), Inches(10.333), Inches(0.5),
     "南京中医药大学 · 软件工程综合实践课程设计",
     size=Pt(15), color=RGBColor(0xDD, 0xCC, 0xBB), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: TOC (restructured to 5 chapters)
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "目  录", "CONTENTS")

chapters = [
    ("01", "项目概述", "背景、目标与系统总体需求"),
    ("02", "功能需求", "核心业务功能与查询统计功能"),
    ("03", "用例分析", "用例图概览与核心用例详细描述"),
    ("04", "数据与架构", "数据库设计 + 4+1架构视图"),
    ("05", "系统特性与总结", "数据安全、设计特点与后续展望"),
]

for i, (num, title, desc) in enumerate(chapters):
    y = Inches(1.8) + Inches(1.0) * i

    # Number circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y + Inches(0.05),
                           Inches(0.55), Inches(0.55))
    c.fill.solid()
    c.fill.fore_color.rgb = DEEP_RED
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = FONT_BODY

    # Title
    tbox(s, Inches(2.35), y + Inches(0.05), Inches(3.5), Inches(0.45),
         title, size=Pt(22), color=DEEP_RED, bold=True, font=FONT_TITLE)

    # Description
    tbox(s, Inches(2.35), y + Inches(0.5), Inches(6), Inches(0.35),
         desc, size=Pt(12), color=MID_GRAY)

    # Separator
    if i < 4:
        gold_line(s, Inches(2.35), y + Inches(0.92), Inches(8.5))

footer(s, 2)


# ═══════════════════════════════════════════════════════════════
# CHAPTER 1: 项目概述 (Slides 3-5)
# ═══════════════════════════════════════════════════════════════

# -- Slide 3: Chapter 1 divider --
chapter_divider("一", "项目概述", "Project Overview",
                "本章介绍项目背景、目标用户、技术路线\n以及系统的六大功能模块总览")

# -- Slide 4: 项目背景与目标 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "项目背景与目标", "BACKGROUND & OBJECTIVES")

# Left: Background
add_rect(s, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), LIGHT_GRAY)
tbox(s, Inches(1.1), Inches(1.75), Inches(5.1), Inches(0.45),
     "▎项目背景", size=Pt(20), color=DEEP_RED, bold=True, font=FONT_TITLE)

bg_items = [
    "中医药是中华民族的瑰宝，药材种类繁多、配伍关系复杂",
    "传统纸质查阅效率低，难以快速准确判断配伍禁忌",
    "中医药信息化是国家中医药发展战略的重要方向",
    "中医药教学和临床实践中急需便捷的数字化辅助工具",
]
for i, item in enumerate(bg_items):
    tbox(s, Inches(1.2), Inches(2.45) + Inches(0.55) * i, Inches(5.0), Inches(0.5),
         f"• {item}", size=Pt(13), color=DARK)

# Right: Goals
add_rect(s, Inches(6.9), Inches(1.6), Inches(5.7), Inches(5.1), LIGHT_GRAY)
tbox(s, Inches(7.2), Inches(1.75), Inches(5.1), Inches(0.45),
     "▎核心目标与定位", size=Pt(20), color=DEEP_RED, bold=True, font=FONT_TITLE)

goals = [
    ("目标用户", "中医药学生、临床中医师、中药药剂师"),
    ("核心定位", "桌面端中医药知识与配伍禁忌辅助工具"),
    ("技术路线", "Python + tkinter + SQLite\n纯内置库实现，零第三方依赖"),
    ("设计理念", "开箱即用、单机运行、操作简洁直观"),
]
for i, (label, desc) in enumerate(goals):
    y = Inches(2.4) + Inches(0.9) * i
    tbox(s, Inches(7.4), y, Inches(5), Inches(0.3),
         label, size=Pt(13), color=GOLD, bold=True)
    tbox(s, Inches(7.4), y + Inches(0.3), Inches(5), Inches(0.55),
         desc, size=Pt(12), color=DARK)

footer(s, 4)

# -- Slide 5: 系统总体需求 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "系统总体需求 — 六大功能模块", "SYSTEM OVERVIEW")

modules = [
    ("药材信息管理", "10字段药材信息的录入、\n修改、删除、浏览"),
    ("多维度药材查询", "名称/拼音/别名/功效/\n分类五维度模糊查询"),
    ("配伍禁忌检查", "多味药材自动配对，\n匹配禁忌规则输出报告"),
    ("禁忌规则管理", "十八反/十九畏/其它类型\n规则的增删改查维护"),
    ("数据统计展示", "自动汇总统计药材总数、\n分类分布、禁忌类型分布"),
    ("数据导入导出", "JSON全量导出导入，\n数据库一键重置恢复"),
]

card_w = Inches(3.65)
card_h = Inches(2.4)
gap_x = Inches(0.22)
gap_y = Inches(0.18)
start_x = Inches(0.8)
start_y = Inches(1.5)

for i, (title, desc) in enumerate(modules):
    col, row = i % 3, i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    add_rect(s, x, y, card_w, card_h, LIGHT_GRAY)
    add_rect(s, x, y, card_w, Pt(4), DEEP_RED)

    # Number badge
    num_badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2),
                                   Inches(0.4), Inches(0.4))
    num_badge.fill.solid()
    num_badge.fill.fore_color.rgb = DEEP_RED
    num_badge.line.fill.background()
    tf = num_badge.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = FONT_BODY

    tbox(s, x + Inches(0.75), y + Inches(0.18), card_w - Inches(1.0), Inches(0.45),
         title, size=Pt(18), color=DEEP_RED, bold=True, font=FONT_TITLE)
    tbox(s, x + Inches(0.3), y + Inches(0.85), card_w - Inches(0.6), Inches(1.3),
         desc, size=Pt(12), color=DARK)

footer(s, 5)


# ═══════════════════════════════════════════════════════════════
# CHAPTER 2: 功能需求 (Slides 6-8)
# ═══════════════════════════════════════════════════════════════

# -- Slide 6: Chapter 2 divider --
chapter_divider("二", "功能需求", "Functional Requirements",
                "本章详细说明系统的核心业务功能\n以及多维查询与数据统计功能")

# -- Slide 7: 核心业务功能 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "核心业务功能", "CORE BUSINESS FUNCTIONS")

# Three columns
cols = [
    ("药材信息管理", DEEP_RED, [
        "10个字段：名称★、拼音、别名、性味、归经、功效、功效分类、用法用量、来源、注意事项",
        "名称必填，设置UNIQUE约束，其余字段选填",
        "支持录入、修改、删除、详情查看全流程操作",
        "数据变更后自动刷新列表，表单一键清空",
    ]),
    ("配伍规则管理", DEEP_RED, [
        "规则组成：药材A + 药材B + 禁忌类型 + 说明",
        "禁忌类型分三类：十八反 / 十九畏 / 其它",
        "药材对自动按字母序排序存储，防止正反重复",
        "复合UNIQUE约束 (herb_a, herb_b) 确保数据唯一",
    ]),
    ("配伍禁忌检查 ★", GOLD, [
        "用户从列表中勾选 ≥2 味药材，触发自动检查",
        "双重循环 C(n,2) 枚举所有可能的药材配对组合",
        "配对名称排序后与禁忌规则库逐条比对匹配",
        "输出冲突报告：安全✓ / 发现 N 处配伍禁忌 ✗",
    ]),
]

for i, (title, accent, items) in enumerate(cols):
    x = Inches(0.8) + Inches(4.05) * i
    w = Inches(3.8)

    # Header
    add_rect(s, x, Inches(1.55), w, Inches(0.55), accent)
    tbox(s, x + Inches(0.15), Inches(1.57), w - Inches(0.3), Inches(0.5),
         title, size=Pt(16), color=WHITE, bold=True, font=FONT_TITLE)

    # Items
    for j, item in enumerate(items):
        y = Inches(2.35) + Inches(0.65) * j
        # Bullet indicator
        add_rect(s, x + Inches(0.15), y + Inches(0.15), Pt(4), Pt(4), GOLD if i == 2 else DEEP_RED)
        tbox(s, x + Inches(0.3), y, w - Inches(0.5), Inches(0.6),
             item, size=Pt(11), color=DARK)

footer(s, 7)

# -- Slide 8: 查询与统计功能 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "查询与统计功能", "QUERY & STATISTICS")

# Left: query
tbox(s, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.45),
     "▎5 维模糊查询", size=Pt(20), color=DEEP_RED, bold=True, font=FONT_TITLE)

queries = [
    ("按名称", "LIKE %keyword%，匹配中文名称"),
    ("按拼音", "拼音片段匹配，如 \"gan\" → \"甘草\""),
    ("按功效", "输入\"清热\"、\"补气\"等功效关键词"),
    ("按分类", "按功效分类筛选，如\"补气药\"、\"清热药\""),
    ("按别名", "支持别名/俗称关键词检索"),
]
for i, (field, desc) in enumerate(queries):
    y = Inches(2.2) + Inches(0.55) * i
    tbox(s, Inches(1.1), y, Inches(1.8), Inches(0.4),
         f"▸ {field}", size=Pt(14), color=DEEP_RED, bold=True)
    tbox(s, Inches(2.9), y + Pt(2), Inches(3.8), Inches(0.4),
         desc, size=Pt(11), color=DARK)

# Right: stats
tbox(s, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.45),
     "▎4 类数据统计", size=Pt(20), color=DEEP_RED, bold=True, font=FONT_TITLE)

stats = [
    ("药材总数", "COUNT herbs 表记录总数"),
    ("禁忌规则总数", "COUNT incompatibilities 表记录数"),
    ("功效分类分布", "GROUP BY category，按数量降序排列"),
    ("禁忌类型分布", "GROUP BY rule_type，分层统计"),
]
for i, (label, desc) in enumerate(stats):
    y = Inches(2.2) + Inches(0.7) * i
    add_rect(s, Inches(7.3), y + Inches(0.08), Pt(4), Inches(0.35), GOLD)
    tbox(s, Inches(7.5), y, Inches(4.8), Inches(0.3),
         label, size=Pt(14), color=DEEP_RED, bold=True)
    tbox(s, Inches(7.5), y + Inches(0.3), Inches(4.8), Inches(0.3),
         desc, size=Pt(11), color=MID_GRAY)

# Bottom note
add_rect(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.8), LIGHT_GRAY)
tbox(s, Inches(1.1), Inches(5.8), Inches(11.2), Inches(0.6),
     "辅助特性：查询输入框支持回车键触发查询 | 数据变更后自动刷新全部列表 | 表单一键清空 | 键盘快捷操作",
     size=Pt(12), color=DARK)

footer(s, 8)


# ═══════════════════════════════════════════════════════════════
# CHAPTER 3: 用例分析 (Slides 9-11)
# ═══════════════════════════════════════════════════════════════

# -- Slide 9: Chapter 3 divider --
chapter_divider("三", "用例分析", "Use Case Analysis",
                "本章展示系统顶层用例图\n并深入剖析核心用例的详细流程")

# -- Slide 10: 用例分析概览 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "用例分析概览", "USE CASE OVERVIEW")

# System boundary background
add_rect(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.6),
         fill_color=LIGHT_GRAY)

# System title
tbox(s, Inches(3.0), Inches(1.4), Inches(7.3), Inches(0.4),
     "中医药常用药材查询与配伍禁忌系统",
     size=Pt(14), color=DEEP_RED, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)

# 6 use cases in 3x2 grid
ucs = [
    ("药材信息管理", "增 · 删 · 改 · 查"),
    ("多维度药材查询", "5维模糊检索"),
    ("配伍禁忌检查", "自动配对 · 冲突报告"),
    ("禁忌规则管理", "规则增删改查"),
    ("数据统计展示", "分类分布 · 自动汇总"),
    ("数据导入导出", "JSON · 重置恢复"),
]

for i, (name, desc) in enumerate(ucs):
    col, row = i % 3, i // 3
    x = Inches(1.7) + Inches(3.5) * col
    y = Inches(2.1) + Inches(1.9) * row

    oval = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(2.8), Inches(1.2))
    oval.fill.solid()
    oval.fill.fore_color.rgb = LIGHT_GRAY
    oval.line.color.rgb = DEEP_RED
    oval.line.width = Pt(1.5)
    tf = oval.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(15)
    p.runs[0].font.color.rgb = DEEP_RED
    p.runs[0].font.bold = True
    p.runs[0].font.name = FONT_BODY
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(10)
    p2.runs[0].font.color.rgb = MID_GRAY
    p2.runs[0].font.name = FONT_BODY

# User actor
ub = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(6.3), Inches(2.9), Inches(0.45))
ub.fill.solid()
ub.fill.fore_color.rgb = DEEP_RED
ub.line.fill.background()
tf = ub.text_frame
p = tf.paragraphs[0]
p.text = "用  户（中医药从业者）"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(13)
p.runs[0].font.color.rgb = WHITE
p.runs[0].font.bold = True
p.runs[0].font.name = FONT_BODY

# Core UC note
tbox(s, Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.3),
     "核心用例：UC-01 药材信息录入  |  UC-02 配伍禁忌自动检查  |  UC-03 JSON数据导出与导入",
     size=Pt(11), color=MID_GRAY, align=PP_ALIGN.CENTER)
footer(s, 10)

# -- Slide 11: UC-02 核心用例详解 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "核心用例详解 — UC-02 配伍禁忌自动检查", "USE CASE DETAIL")

# Flow steps
steps = [
    ("① 选择药材", "用户在可选药材列表\n中浏览全部药材名称"),
    ("② 添加待检", "双击或点击按钮将药材\n加入待检查列表（≥2味）"),
    ("③ 执行检查", "点击\"执行配伍禁忌\n检查\"触发自动匹配"),
    ("④ 枚举配对", "双重循环C(n,2)枚举\n排序后与规则库比对"),
    ("⑤ 输出报告", "无冲突→安全提示\n有冲突→逐条列出详情"),
]

for i, (name, desc) in enumerate(steps):
    x = Inches(0.4) + Inches(2.55) * i

    # Step circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.6),
                           Inches(0.55), Inches(0.55))
    c.fill.solid()
    c.fill.fore_color.rgb = GOLD if i == 3 else DEEP_RED
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(18)
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = FONT_BODY

    tbox(s, x + Inches(0.05), Inches(2.3), Inches(2.45), Inches(0.4),
         name, size=Pt(14), color=DEEP_RED, bold=True, align=PP_ALIGN.CENTER)
    tbox(s, x + Inches(0.05), Inches(2.7), Inches(2.45), Inches(1.0),
         desc, size=Pt(11), color=DARK, align=PP_ALIGN.CENTER)

    if i < 4:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   x + Inches(2.35), Inches(1.75),
                                   Inches(0.25), Inches(0.22))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

# Design key points
add_rect(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.6), LIGHT_GRAY)
tbox(s, Inches(1.1), Inches(4.1), Inches(11.2), Inches(0.4),
     "▎关键设计要点", size=Pt(18), color=DEEP_RED, bold=True, font=FONT_TITLE)

notes = [
    "名称排序匹配：输入药材名称先strip()去空格再排序，规则存储时herb_a/herb_b也按字母序排列，保证正反序输入均能正确匹配",
    "双重循环算法：O(n²)枚举所有药材配对，实际场景中通常<50味药材，性能完全满足需求",
    "规则在Python层匹配：不在数据库层做JOIN，而是在Python层获取所有规则后做名称比对，逻辑更灵活",
    "结果分级展示：无冲突时绿色提示\"未发现配伍禁忌\"，有冲突时红色逐条列出冲突药材对、禁忌类型及说明",
]
for i, note in enumerate(notes):
    tbox(s, Inches(1.3), Inches(4.6) + Inches(0.43) * i, Inches(10.9), Inches(0.4),
         f"• {note}", size=Pt(12), color=DARK)

footer(s, 11)


# ═══════════════════════════════════════════════════════════════
# CHAPTER 4: 数据与架构 (Slides 12-19)
# ═══════════════════════════════════════════════════════════════

# -- Slide 12: Chapter 4 divider --
chapter_divider("四", "数据与架构", "Data & Architecture",
                "本章展示数据库表结构设计与数据字典\n以及系统的4+1架构视图")

# -- Slide 13: 数据库设计(一) — herbs表 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "数据库设计（一）— herbs 药材信息表", "DATABASE DESIGN")

# Table data
herb_cols = [
    ("序号", "字段名", "类型", "约束", "说明"),
    ("1", "id", "INTEGER", "PK AUTOINCREMENT", "药材唯一标识"),
    ("2", "name", "TEXT", "NOT NULL UNIQUE", "中药名称（如\"甘草\"）★"),
    ("3", "pinyin", "TEXT", "DEFAULT ''", "汉语拼音（如\"gancao\"）"),
    ("4", "alias", "TEXT", "DEFAULT ''", "别名/俗称（如\"国老\"）"),
    ("5", "xingwei", "TEXT", "DEFAULT ''", "性味（如\"甘，平\"）"),
    ("6", "guijing", "TEXT", "DEFAULT ''", "归经（如\"心、肺、脾\"）"),
    ("7", "gongxiao", "TEXT", "DEFAULT ''", "功效描述"),
    ("8", "category", "TEXT", "DEFAULT ''", "功效分类（如\"补气药\"）"),
    ("9", "yongfa_yongliang", "TEXT", "DEFAULT ''", "用法用量"),
    ("10", "laiyuan", "TEXT", "DEFAULT ''", "来源/基原"),
    ("11", "zhuyi", "TEXT", "DEFAULT ''", "注意事项/禁忌"),
]

# Calculate column widths - wider for long field names
cw = [Inches(0.7), Inches(2.2), Inches(1.3), Inches(2.8), Inches(5.3)]
tx = Inches(0.65)
ty = Inches(1.55)
rh = Inches(0.38)

# Header row
for j, h in enumerate(herb_cols[0]):
    x = tx + sum(cw[:j])
    add_rect(s, x, ty, cw[j], rh, DEEP_RED)
    tbox(s, x + Inches(0.04), ty + Pt(3), cw[j] - Inches(0.08), rh,
         h, size=Pt(11), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(herb_cols[1:]):
    for j, cell in enumerate(row):
        x = tx + sum(cw[:j])
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rect(s, x, ty + rh * (i + 1), cw[j], rh, bg)
        c = DEEP_RED if j == 1 else DARK
        b = j == 1
        align = PP_ALIGN.CENTER if j in [0, 2] else PP_ALIGN.LEFT
        tbox(s, x + Inches(0.06), ty + rh * (i + 1) + Pt(3),
             cw[j] - Inches(0.12), rh,
             cell, size=Pt(11), color=c, bold=b, align=align)

# Notes below table
tbox(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7),
     "说明：★ 为必填字段（NOT NULL + UNIQUE约束）。功效分类涵盖解表药、清热药、补虚药（补气/补血/补阴/补阳）等。"
     "用法用量、来源、注意事项为多行文本字段（tkinter Text height=2）。",
     size=Pt(11), color=MID_GRAY)

footer(s, 13)

# -- Slide 14: 数据库设计(二) — incompatibilities表 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "数据库设计（二）— incompatibilities 配伍禁忌规则表", "DATABASE DESIGN")

inc_cols = [
    ("序号", "字段名", "类型", "约束", "说明"),
    ("1", "id", "INTEGER", "PK AUTOINCREMENT", "规则唯一标识"),
    ("2", "herb_a", "TEXT", "NOT NULL", "药材A（按字母序存储）"),
    ("3", "herb_b", "TEXT", "NOT NULL", "药材B（按字母序存储）"),
    ("4", "rule_type", "TEXT", "DEFAULT ''", "十八反 / 十九畏 / 其它"),
    ("5", "description", "TEXT", "DEFAULT ''", "禁忌说明（如\"甘草反甘遂\"）"),
    ("—", "(herb_a, herb_b)", "—", "UNIQUE", "复合唯一约束，防止重复"),
]

cw = [Inches(0.7), Inches(2.8), Inches(1.3), Inches(2.8), Inches(5.7)]
tx = Inches(0.65)
ty = Inches(1.55)
rh = Inches(0.44)

for j, h in enumerate(inc_cols[0]):
    x = tx + sum(cw[:j])
    add_rect(s, x, ty, cw[j], rh, DEEP_RED)
    tbox(s, x + Inches(0.04), ty + Pt(4), cw[j] - Inches(0.08), rh,
         h, size=Pt(11), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(inc_cols[1:]):
    for j, cell in enumerate(row):
        x = tx + sum(cw[:j])
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rect(s, x, ty + rh * (i + 1), cw[j], rh, bg)
        c = DEEP_RED if j == 1 else DARK
        b = j == 1
        align = PP_ALIGN.CENTER if j in [0, 2] else PP_ALIGN.LEFT
        tbox(s, x + Inches(0.06), ty + rh * (i + 1) + Pt(4),
             cw[j] - Inches(0.12), rh,
             cell, size=Pt(11), color=c, bold=b, align=align)

# Design note
tbox(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.3),
     "▎配伍规则存储约定", size=Pt(16), color=DEEP_RED, bold=True, font=FONT_TITLE)

notes = [
    "INSERT时自动将herb_a与herb_b按字母序排序后存储，确保（甘草, 甘遂）与（甘遂, 甘草）识别为同一条规则",
    "配伍检查时同样对输入药材名称排序后比对，保证正反序输入均能正确匹配",
    "禁忌类型参照中医经典分类：十八反（乌头反半夏等6组）、十九畏（硫黄畏朴硝等9组）、其它自定义类型",
    "描述字段存储具体禁忌说明，如\"甘草反甘遂，两药合用可增强毒性\"",
]
for i, note in enumerate(notes):
    tbox(s, Inches(1.1), Inches(5.4) + Inches(0.38) * i, Inches(11), Inches(0.35),
         f"• {note}", size=Pt(12), color=DARK)

footer(s, 14)

# -- Slides 15-19: Architecture Views --
arch_slides = [
    ("逻辑视图", "逻辑视图.png",
     "从功能分解角度描述系统六大模块及其关系，"
     "展示药材管理、查询、配伍检查、规则管理、统计、数据操作模块的职责划分与依赖。"),
    ("开发视图", "开发视图.png",
     "从代码组织角度展示系统的分层结构："
     "表现层（ui.py）、业务逻辑层（database.py）、数据层（SQLite），"
     "以及种子数据（data.py）和程序入口（main.py）。"),
    ("物理视图", "物理视图.png",
     "从部署角度描述系统的物理拓扑："
     "单机桌面应用，Python运行时直接访问本地SQLite数据库文件（tcm.db），无需网络通信。"),
    ("场景视图", "场景视图.png",
     "从运行时刻角度展示关键用例的对象交互："
     "用户操作→GUI事件→DBManager方法调用→SQLite执行→结果返回→界面刷新。"),
    ("核心业务流程时序图", "核心业务流程时序图.png",
     "以时序图形式展示\"配伍禁忌检查\"的完整交互流程："
     "用户选药→UI层收集名称→DBManager查询规则→Python层配对匹配→结果回传展示。"),
]

for idx, (title, filename, desc) in enumerate(arch_slides):
    s = add_blank_slide()
    add_bg(s, WHITE)
    section_title(s, f"系统架构 — {title}", "ARCHITECTURE VIEW")

    # Pre-calculate image size for centering
    img_path = os.path.join(IMG_DIR, filename)
    img_max_w = Inches(10.5)
    img_max_h = Inches(4.6)
    img_top = Inches(1.5)

    if os.path.exists(img_path):
        from PIL import Image
        with Image.open(img_path) as img:
            pw, ph = img.size
        aspect = pw / ph if ph > 0 else 1
        if img_max_w / aspect <= img_max_h:
            actual_w = int(img_max_w)
            actual_h = int(img_max_w / aspect)
        else:
            actual_h = int(img_max_h)
            actual_w = int(img_max_h * aspect)
        # Center horizontally
        img_left = (W - actual_w) // 2
        s.shapes.add_picture(img_path, img_left, img_top, actual_w, actual_h)

        desc_top = img_top + actual_h + Inches(0.2)
    else:
        actual_h = None
        desc_top = Inches(1.5)

    # Description text centered below image
    tbox(s, Inches(1.5), desc_top, Inches(10.333), Inches(0.8),
         desc, size=Pt(12), color=DARK, align=PP_ALIGN.CENTER)

    footer(s, 15 + idx)


# ═══════════════════════════════════════════════════════════════
# CHAPTER 5: 系统特性与总结 (Slides 20-22)
# ═══════════════════════════════════════════════════════════════

# -- Slide 20: Chapter 5 divider --
chapter_divider("五", "系统特性与总结", "Features & Summary",
                "本章介绍系统的设计特点与数据安全保障\n并对需求分析工作进行总结与展望")

# -- Slide 21: 系统特点与数据安全 --
s = add_blank_slide()
add_bg(s, WHITE)
section_title(s, "系统特点与数据安全", "FEATURES & DATA SAFETY")

features = [
    ("零依赖 · 开箱即用",
     "纯Python内置库（sqlite3, tkinter, json）\n无需pip安装，拷贝文件夹即运行"),
    ("测试隔离 · 安全可靠",
     "test.py使用tempfile临时数据库，\n测试过程完全不影响生产数据 tcm.db"),
    ("操作确认 · 防误操作",
     "删除药材、删除规则、重置数据库等\n不可逆操作均弹出确认对话框再执行"),
    ("数据约束 · 保证一致",
     "药材名称 UNIQUE + 规则对复合 UNIQUE，\n从数据库层面杜绝重复录入"),
    ("备份恢复 · 数据无忧",
     "支持JSON全量导出/导入，\n用户可定期备份，支持数据迁移恢复"),
    ("首次初始化 · 即装即用",
     "首次运行自动检测空数据库，\n自动导入25种种子药材 + 17条禁忌规则"),
]

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.1) * col
    y = Inches(1.65) + Inches(1.7) * row

    # Number circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
    c.fill.solid()
    c.fill.fore_color.rgb = DEEP_RED
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.size = Pt(16)
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.name = FONT_BODY

    tbox(s, x + Inches(0.65), y + Pt(4), Inches(5), Inches(0.4),
         title, size=Pt(18), color=DEEP_RED, bold=True, font=FONT_TITLE)
    tbox(s, x + Inches(0.65), y + Inches(0.55), Inches(5), Inches(0.9),
         desc, size=Pt(12), color=DARK)

footer(s, 21)

# -- Slide 22: Summary & Outlook --
s = add_blank_slide()
add_bg(s, DEEP_RED)
add_rect(s, Inches(0), Inches(0), W, Pt(8), GOLD)
add_rect(s, Inches(0), H - Pt(8), W, Pt(8), GOLD)

tbox(s, Inches(1.5), Inches(0.9), Inches(10.333), Inches(0.8),
     "总结与展望", size=Pt(36), color=WHITE, bold=True,
     font=FONT_TITLE, align=PP_ALIGN.CENTER)
gold_line(s, Inches(5.2), Inches(1.7), Inches(2.933))

summary_items = [
    "系统定位：面向中医药从业者的桌面端辅助工具，纯Python零依赖实现",
    "功能需求：明确6大功能模块，覆盖药材管理、查询、配伍检查等全流程",
    "用例分析：梳理3个核心用例，含详细的基本流程与异常处理机制",
    "数据设计：2张SQLite数据表，11+5字段，约束设计完善",
    "架构视图：完成4+1架构视图，清晰展示系统多维度结构",
]

for i, item in enumerate(summary_items):
    tbox(s, Inches(1.8), Inches(2.1) + Inches(0.55) * i, Inches(9.733), Inches(0.5),
         f"✓ {item}", size=Pt(15), color=WHITE)

# Next steps
tbox(s, Inches(1.5), Inches(5.1), Inches(10.333), Inches(0.5),
     "后续工作", size=Pt(22), color=GOLD, bold=True,
     font=FONT_TITLE, align=PP_ALIGN.CENTER)
tbox(s, Inches(1.5), Inches(5.6), Inches(10.333), Inches(0.5),
     "系统设计  →  编码实现  →  测试验证  →  报告撰写",
     size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)

tbox(s, Inches(1.5), Inches(6.3), Inches(10.333), Inches(0.6),
     "感谢聆听 · 敬请指导", size=Pt(26), color=GOLD,
     font=FONT_TITLE, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"PPT saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
